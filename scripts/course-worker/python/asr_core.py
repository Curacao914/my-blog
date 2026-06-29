#!/usr/bin/env python3
from __future__ import annotations

import getpass
import hashlib
import json
import math
import os
import re
import shutil
import socket
import ssl
import subprocess
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib import error, request

import boto3
from botocore.config import Config
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

ROOT = Path(__file__).resolve().parents[1]
ENV_PATH = ROOT / ".env.local"
LOG_DIR = ROOT / "logs"
DIAG_DIR = ROOT / "diagnostics"
OUTPUT_ROOT = ROOT / "outputs"
PRIVATE_DIR = ROOT / ".private"
INBOX_DIR = ROOT / "materials" / "inbox"

DEFAULT_R2_ENDPOINT = ""
DEFAULT_R2_BUCKET = ""
DASHSCOPE_SUBMIT = "https://dashscope.aliyuncs.com/api/v1/services/audio/asr/transcription"
DASHSCOPE_TASK = "https://dashscope.aliyuncs.com/api/v1/tasks/"
MODEL = "paraformer-v2"
PRICE_PER_HOUR_CNY = 0.288
HOME = str(Path.home())
SUPPORTED_MATERIALS = {".pptx", ".pdf", ".docx", ".md", ".txt"}
RETRYABLE_HTTP = {408, 425, 429, 500, 502, 503, 504}
URL_RE = re.compile(r"https?://[^\s'\"<>]+", re.IGNORECASE)
JWT_RE = re.compile(r"eyJ[A-Za-z0-9_-]{10,}\.[A-Za-z0-9_-]{10,}(?:\.[A-Za-z0-9_-]{10,})?")


def sanitize_text(text: str) -> str:
    value = text or ""
    value = value.replace(HOME, "<HOME>")
    value = URL_RE.sub("<REDACTED_URL>", value)
    value = JWT_RE.sub("<REDACTED_JWT>", value)
    value = re.sub(
        r"(?i)(api[_-]?key|secret[_-]?access[_-]?key|access[_-]?key[_-]?id)\s*[:=]\s*[^\s,}]+",
        r"\1=<REDACTED>",
        value,
    )
    value = re.sub(r"(?i)(authorization\s*:\s*bearer\s+)[^\s]+", r"\1<REDACTED>", value)
    return value


def sanitize_json(value: Any) -> Any:
    if isinstance(value, dict):
        clean: dict[str, Any] = {}
        for key, item in value.items():
            lowered = str(key).lower()
            if lowered in {
                "file_url", "file_urls", "transcription_url", "url",
                "authorization", "api_key", "secret_access_key", "access_key_id",
                "task_id",
            }:
                clean[key] = "<REDACTED_URL>" if "url" in lowered else "<REDACTED>"
            else:
                clean[key] = sanitize_json(item)
        return clean
    if isinstance(value, list):
        return [sanitize_json(item) for item in value]
    if isinstance(value, str):
        return sanitize_text(value)
    return value


def safe_slug(value: str) -> str:
    value = re.sub(r"\s+", "-", value.strip())
    value = re.sub(r"[^\w\u4e00-\u9fff-]+", "-", value)
    return re.sub(r"-{2,}", "-", value).strip("-")[:100] or "lesson"


def format_timestamp(milliseconds: int) -> str:
    total = max(0, int(milliseconds // 1000))
    hours, remainder = divmod(total, 3600)
    minutes, seconds = divmod(remainder, 60)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d}"


def load_env_file(path: Path) -> dict[str, str]:
    values: dict[str, str] = {}
    if not path.exists():
        return values
    for line in path.read_text(encoding="utf-8").splitlines():
        stripped = line.strip()
        if stripped and not stripped.startswith("#") and "=" in stripped:
            key, raw = stripped.split("=", 1)
            values[key.strip()] = raw.strip()
    return values


def previous_v004_dirs() -> list[Path]:
    return sorted(
        [
            p for p in ROOT.parent.glob("law-tech-course-pipeline-debug-v004-*")
            if p.is_dir() and p.resolve() != ROOT.resolve()
        ],
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    )


def maybe_import_previous_env() -> None:
    if ENV_PATH.exists():
        return
    for previous in previous_v004_dirs():
        candidate = previous / ".env.local"
        if candidate.exists():
            print(f"发现旧 V004 密钥配置：{previous.name}")
            if (input("复用？[Y/n]: ").strip().lower() or "y") in {"y", "yes"}:
                shutil.copy2(candidate, ENV_PATH)
                os.chmod(ENV_PATH, 0o600)
                print("✓ 已复用密钥配置。")
            return
    for candidate in sorted(
        ROOT.parent.glob("law-tech-course-pipeline-debug-v003-paraformer*/.env.local"),
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    ):
        print(f"发现 V003 密钥配置：{candidate.parent.name}")
        if (input("复用？[Y/n]: ").strip().lower() or "y") in {"y", "yes"}:
            shutil.copy2(candidate, ENV_PATH)
            os.chmod(ENV_PATH, 0o600)
            print("✓ 已复用密钥配置。")
        return


def configure_secrets() -> dict[str, str]:
    maybe_import_previous_env()
    values = load_env_file(ENV_PATH)
    required = ("DASHSCOPE_API_KEY", "R2_ACCESS_KEY_ID", "R2_SECRET_ACCESS_KEY")
    if all(values.get(key) for key in required):
        return values
    print("输入不会回显，密钥只保存在 .env.local。")
    values = {
        "DASHSCOPE_API_KEY": getpass.getpass("阿里云百炼 DASHSCOPE_API_KEY: ").strip(),
        "R2_ACCESS_KEY_ID": getpass.getpass("R2 Access Key ID: ").strip(),
        "R2_SECRET_ACCESS_KEY": getpass.getpass("R2 Secret Access Key: ").strip(),
        "R2_ENDPOINT": DEFAULT_R2_ENDPOINT,
        "R2_BUCKET": DEFAULT_R2_BUCKET,
    }
    if not all(values.get(key) for key in required):
        raise RuntimeError("三个密钥均不能为空")
    ENV_PATH.write_text(
        "\n".join(f"{key}={value}" for key, value in values.items()) + "\n",
        encoding="utf-8",
    )
    os.chmod(ENV_PATH, 0o600)
    return values


def is_retryable_exception(exc: BaseException) -> bool:
    if isinstance(exc, error.HTTPError):
        return exc.code in RETRYABLE_HTTP
    return isinstance(exc, (
        error.URLError, ssl.SSLError, socket.timeout, TimeoutError,
        ConnectionResetError, ConnectionAbortedError, BrokenPipeError,
    ))


def urlopen_json(req: request.Request, timeout: int = 90, attempts: int = 6, log=None) -> dict[str, Any]:
    last_error: BaseException | None = None
    for attempt in range(1, attempts + 1):
        try:
            with request.urlopen(req, timeout=timeout) as response:
                return json.loads(response.read().decode("utf-8"))
        except error.HTTPError as exc:
            body = exc.read().decode("utf-8", errors="replace")
            last_error = RuntimeError(f"HTTP {exc.code}: {sanitize_text(body)}")
            retryable = exc.code in RETRYABLE_HTTP
        except BaseException as exc:
            last_error = exc
            retryable = is_retryable_exception(exc)
        if not retryable or attempt == attempts:
            break
        wait = min(30, 2 ** (attempt - 1))
        if log:
            log(f"NETWORK_RETRY attempt={attempt} wait={wait}s error={last_error}")
        print(f"      网络连接短暂中断，{wait} 秒后重试（{attempt}/{attempts}）")
        time.sleep(wait)
    raise RuntimeError("网络请求失败：" + sanitize_text(str(last_error))) from last_error


def post_json(url: str, payload: dict[str, Any], headers: dict[str, str], log=None) -> dict[str, Any]:
    return urlopen_json(request.Request(
        url,
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers=headers,
        method="POST",
    ), log=log)


def get_json(url: str, headers: dict[str, str] | None = None, log=None) -> dict[str, Any]:
    return urlopen_json(request.Request(url, headers=headers or {}, method="GET"), log=log)


def run_command(args: list[str], timeout: int, log) -> subprocess.CompletedProcess[str]:
    safe = [
        "<REDACTED_LOCAL_PATH>" if item.startswith(HOME) else
        "<REDACTED_URL>" if item.startswith(("http://", "https://")) else item
        for item in args
    ]
    log("RUN " + " ".join(safe))
    started = time.monotonic()
    result = subprocess.run(
        args, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        text=True, timeout=timeout, check=False,
    )
    log(f"EXIT {result.returncode} in {time.monotonic() - started:.2f}s")
    if result.stderr:
        log("STDERR " + sanitize_text(result.stderr[-4000:]))
    return result


def find_recent_files(extensions: set[str]) -> list[Path]:
    found: dict[str, Path] = {}
    INBOX_DIR.mkdir(parents=True, exist_ok=True)
    for location in (Path.home() / "Downloads", INBOX_DIR):
        if location.exists():
            for path in location.iterdir():
                if path.is_file() and path.suffix.lower() in extensions:
                    found[str(path.resolve())] = path.resolve()
    return sorted(found.values(), key=lambda p: p.stat().st_mtime, reverse=True)


def choose_one(title: str, extensions: set[str]) -> Path:
    files = find_recent_files(extensions)
    if not files:
        raise RuntimeError(f"没有找到{title}")
    print(f"\n{title}：")
    for i, path in enumerate(files[:12], 1):
        print(f"  {i}. {path.name}")
    raw = input("选择 [1]: ").strip() or "1"
    return files[int(raw) - 1]


def choose_materials() -> list[Path]:
    files = find_recent_files(SUPPORTED_MATERIALS)
    if not files:
        raise RuntimeError("没有找到课程 PPTX/PDF/DOCX/MD/TXT")
    print("\n最近的课程材料：")
    for i, path in enumerate(files[:15], 1):
        print(f"  {i}. {path.name}")
    raw = input("选择一个或多个，例如 1 或 1,2 [1]: ").strip() or "1"
    return [files[int(part.strip()) - 1] for part in raw.split(",")]


def extract_material(path: Path) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    sections: list[dict[str, Any]] = []
    ocr: list[dict[str, Any]] = []
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        deck = Presentation(str(path))
        for i, slide in enumerate(deck.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = "\n".join(
                        p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                    )
                    if value:
                        texts.append(value)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            text = "\n".join(texts).strip()
            sections.append({"label": f"第 {i} 页", "text": text})
            if len(re.sub(r"\s+", "", text)) < 12:
                ocr.append({"source": path.name, "kind": "pptx-slide", "index": i})
    elif suffix == ".pdf":
        reader = PdfReader(str(path))
        for i, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            sections.append({"label": f"第 {i} 页", "text": text})
            if len(re.sub(r"\s+", "", text)) < 20:
                ocr.append({"source": path.name, "kind": "pdf-page", "index": i})
    elif suffix == ".docx":
        doc = Document(str(path))
        blocks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    blocks.append(" | ".join(cells))
        sections = [{"label": "全文", "text": "\n".join(blocks)}]
    else:
        sections = [{"label": "全文", "text": path.read_text(encoding="utf-8", errors="replace")}]
    return sections, ocr


def write_materials(output_dir: Path, paths: list[Path]) -> tuple[str, list[dict[str, Any]]]:
    lines = ["# 课程材料文字", ""]
    all_text: list[str] = []
    all_ocr: list[dict[str, Any]] = []
    for path in paths:
        sections, ocr = extract_material(path)
        lines.extend([f"## {path.stem}", ""])
        for section in sections:
            lines.extend([f"### {section['label']}", "", section["text"] or "（未提取到文字）", ""])
            all_text.append(section["text"])
        all_ocr.extend(ocr)
    content = "\n".join(lines).rstrip() + "\n"
    (output_dir / "course-materials.md").write_text(content, encoding="utf-8")
    (output_dir / "ocr-needed.json").write_text(
        json.dumps(all_ocr, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return content, all_ocr


def maybe_import_previous_output(lesson_id: str, output_dir: Path) -> None:
    if output_dir.exists() and any(output_dir.iterdir()):
        return
    for previous in previous_v004_dirs():
        source = previous / "outputs" / lesson_id
        if source.exists():
            print(f"发现旧 V004 输出：{source}")
            if (input("迁移已有检查点？[Y/n]: ").strip().lower() or "y") in {"y", "yes"}:
                shutil.copytree(source, output_dir, dirs_exist_ok=True)
                print("✓ 已迁移旧输出，第 1 段等已完成检查点将直接复用。")
            return


def ffprobe_duration(source: Path, log) -> float:
    result = run_command([
        "ffprobe", "-v", "error", "-show_entries", "format=duration",
        "-of", "default=noprint_wrappers=1:nokey=1", str(source)
    ], 120, log)
    if result.returncode != 0:
        raise RuntimeError("ffprobe 无法读取课程视频")
    return float(result.stdout.strip())


def test_presigned_url(url: str, log) -> None:
    last: BaseException | None = None
    for attempt in range(1, 5):
        try:
            req = request.Request(url, method="GET", headers={"Range": "bytes=0-1023"})
            with request.urlopen(req, timeout=30) as response:
                if getattr(response, "status", 200) not in (200, 206):
                    raise RuntimeError("签名 URL 状态异常")
                response.read(1024)
                return
        except BaseException as exc:
            last = exc
            if attempt == 4 or not is_retryable_exception(exc):
                break
            wait = 2 ** (attempt - 1)
            log(f"PRESIGNED_RETRY attempt={attempt} error={exc}")
            time.sleep(wait)
    raise RuntimeError("R2 签名地址测试失败：" + sanitize_text(str(last)))


def extract_sentences(result: dict[str, Any], offset_ms: int) -> tuple[list[dict[str, Any]], int, int]:
    sentences: list[dict[str, Any]] = []
    original_ms = int((result.get("properties") or {}).get("original_duration_in_milliseconds") or 0)
    speech_ms = 0
    for transcript in result.get("transcripts") or []:
        speech_ms += int(transcript.get("content_duration_in_milliseconds") or 0)
        for sentence in transcript.get("sentences") or []:
            text = str(sentence.get("text") or "").strip()
            if text:
                sentences.append({
                    "begin_time": int(sentence.get("begin_time") or 0) + offset_ms,
                    "end_time": int(sentence.get("end_time") or 0) + offset_ms,
                    "text": text,
                })
    return sentences, original_ms, speech_ms


def task_failure_summary(task_response: dict[str, Any]) -> str:
    output = task_response.get("output") or {}
    parts: list[str] = []
    for key in ("code", "message", "task_status"):
        value = output.get(key)
        if value:
            parts.append(f"{key}={value}")
    for item in output.get("results") or []:
        for key in ("subtask_status", "code", "message"):
            value = item.get(key)
            if value:
                parts.append(f"{key}={value}")
    return "; ".join(dict.fromkeys(parts)) or "未返回具体失败原因"


def safe_delete_r2_object(s3, bucket: str, object_key: str, chunk_index: int, log) -> None:
    try:
        s3.delete_object(Bucket=bucket, Key=object_key)
        log(f"CHUNK {chunk_index} CLEANUP r2_object_deleted=true")
    except Exception as exc:
        log(f"CHUNK {chunk_index} CLEANUP r2_object_deleted=false {exc}")
        raise RuntimeError(f"第 {chunk_index} 段 R2 临时对象删除失败") from exc


def transcribe_chunk(
    source: Path, chunk_index: int, start_seconds: float, duration_seconds: float,
    checkpoint_path: Path, task_path: Path, temp_dir: Path,
    config: dict[str, str], s3, bucket: str, log,
) -> dict[str, Any]:
    if checkpoint_path.exists():
        return json.loads(checkpoint_path.read_text(encoding="utf-8"))

    failure_path = checkpoint_path.with_suffix(".failure.json")
    task_state: dict[str, Any] | None = None

    if task_path.exists():
        loaded = json.loads(task_path.read_text(encoding="utf-8"))
        # R2 task files did not retain the R2 object key and are therefore not resumable.
        if loaded.get("object_key"):
            task_state = loaded
            print(f"   发现第 {chunk_index} 段已提交任务，继续轮询 …{task_state['task_id'][-8:]}")
        else:
            stale_path = task_path.with_suffix(".stale-r2.json")
            task_path.replace(stale_path)
            print(f"   发现旧版不可恢复任务状态，已归档并重新提交第 {chunk_index} 段")

    if task_state is None:
        audio_path = temp_dir / f"chunk-{chunk_index:03d}.mp3"
        object_key = f"course-pipeline-temp/{uuid.uuid4().hex}.mp3"
        uploaded = False
        try:
            result = run_command([
                "ffmpeg", "-nostdin", "-hide_banner", "-loglevel", "error",
                "-ss", f"{start_seconds:.3f}", "-i", str(source),
                "-t", f"{duration_seconds:.3f}", "-vn", "-ac", "1", "-ar", "16000",
                "-c:a", "libmp3lame", "-b:a", "48k", "-y", str(audio_path),
            ], 900, log)
            if result.returncode != 0 or not audio_path.exists() or audio_path.stat().st_size == 0:
                raise RuntimeError(f"第 {chunk_index} 段音频提取失败")

            s3.upload_file(
                str(audio_path),
                bucket,
                object_key,
                ExtraArgs={"ContentType": "audio/mpeg"},
            )
            uploaded = True
            signed_url = s3.generate_presigned_url(
                "get_object",
                Params={"Bucket": bucket, "Key": object_key},
                ExpiresIn=21600,
            )
            test_presigned_url(signed_url, log)

            submitted = post_json(
                DASHSCOPE_SUBMIT,
                {
                    "model": MODEL,
                    "input": {"file_urls": [signed_url]},
                    "parameters": {
                        "channel_id": [0],
                        "language_hints": ["zh", "en"],
                        "disfluency_removal_enabled": False,
                        "timestamp_alignment_enabled": False,
                        "diarization_enabled": False,
                        "special_word_filter": json.dumps({
                            "filter_with_signed": {"word_list": []},
                            "filter_with_empty": {"word_list": []},
                            "system_reserved_filter": False,
                        }, ensure_ascii=False),
                    },
                },
                {
                    "Authorization": f"Bearer {config['DASHSCOPE_API_KEY']}",
                    "Content-Type": "application/json",
                    "X-DashScope-Async": "enable",
                },
                log=log,
            )
            task_id = str((submitted.get("output") or {}).get("task_id") or "")
            if not task_id:
                raise RuntimeError("提交响应没有 task_id")

            # Persist both task id and R2 object key before polling. The object must remain
            # available until DashScope has downloaded and completed the task.
            task_state = {
                "schemaVersion": 2,
                "task_id": task_id,
                "object_key": object_key,
                "bucket": bucket,
                "chunkIndex": chunk_index,
                "startSeconds": start_seconds,
                "durationSecondsRequested": duration_seconds,
                "submittedAt": datetime.now(timezone.utc).isoformat(),
                "signedUrlExpiresSeconds": 21600,
            }
            task_path.write_text(
                json.dumps(task_state, ensure_ascii=False, indent=2) + "\n",
                encoding="utf-8",
            )
            print(f"   第 {chunk_index} 段任务 …{task_id[-8:]} 已提交；临时音频将保留到任务结束")
        except Exception:
            # If submission did not persist a resumable task, clean any uploaded object.
            if uploaded and task_state is None:
                safe_delete_r2_object(s3, bucket, object_key, chunk_index, log)
            raise
        finally:
            if audio_path.exists():
                audio_path.unlink()
                log(f"CHUNK {chunk_index} CLEANUP local_audio_deleted=true")

    assert task_state is not None
    task_id = task_state["task_id"]
    object_key = task_state["object_key"]
    task_bucket = task_state.get("bucket") or bucket
    headers = {
        "Authorization": f"Bearer {config['DASHSCOPE_API_KEY']}",
        "Content-Type": "application/json",
    }

    deadline = time.monotonic() + 40 * 60
    last_status = None
    task_response: dict[str, Any] | None = None

    try:
        while time.monotonic() < deadline:
            time.sleep(4)
            task_response = urlopen_json(
                request.Request(
                    DASHSCOPE_TASK + task_id,
                    headers=headers,
                    data=b"",
                    method="POST",
                ),
                log=log,
            )
            status = str(
                ((task_response.get("output") or {}).get("task_status")) or "UNKNOWN"
            ).upper()
            if status != last_status:
                print(f"      状态：{status}")
                log(f"CHUNK {chunk_index} TASK_STATUS {status}")
                last_status = status

            if status == "SUCCEEDED":
                break

            if status == "FAILED":
                detail = task_failure_summary(task_response)
                failure_payload = {
                    "chunkIndex": chunk_index,
                    "taskIdSuffix": task_id[-8:],
                    "failedAt": datetime.now(timezone.utc).isoformat(),
                    "detail": sanitize_text(detail),
                    "response": sanitize_json(task_response),
                }
                failure_path.write_text(
                    json.dumps(failure_payload, ensure_ascii=False, indent=2) + "\n",
                    encoding="utf-8",
                )
                safe_delete_r2_object(s3, task_bucket, object_key, chunk_index, log)
                task_path.unlink(missing_ok=True)
                raise RuntimeError(f"第 {chunk_index} 段任务失败：{sanitize_text(detail)}")

            if status not in {"PENDING", "RUNNING"}:
                raise RuntimeError(
                    f"第 {chunk_index} 段出现未知任务状态 {status}；task 状态已保留"
                )
        else:
            raise RuntimeError(
                f"第 {chunk_index} 段等待超过 40 分钟；task 与 R2 对象均已保留，可重跑"
            )

        results = ((task_response or {}).get("output") or {}).get("results") or []
        successful = [
            item for item in results if item.get("subtask_status") == "SUCCEEDED"
        ]
        if not successful or not successful[0].get("transcription_url"):
            raise RuntimeError(
                f"第 {chunk_index} 段任务成功但没有可下载结果；task 状态已保留"
            )

        result = get_json(successful[0]["transcription_url"], log=log)
        sentences, original_ms, speech_ms = extract_sentences(
            result,
            int(start_seconds * 1000),
        )
        checkpoint = {
            "chunkIndex": chunk_index,
            "startSeconds": start_seconds,
            "durationSecondsRequested": duration_seconds,
            "taskIdSuffix": task_id[-8:],
            "originalDurationMilliseconds": original_ms,
            "speechDurationMilliseconds": speech_ms,
            "sentences": sentences,
        }
        checkpoint_path.write_text(
            json.dumps(checkpoint, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )

        safe_delete_r2_object(s3, task_bucket, object_key, chunk_index, log)
        task_path.unlink(missing_ok=True)
        failure_path.unlink(missing_ok=True)
        return checkpoint

    except KeyboardInterrupt:
        print(
            f"\n已中断；第 {chunk_index} 段 task 和 R2 临时对象仍保留，"
            "下次运行会继续，不会重新提交。"
        )
        raise

def write_transcript(path: Path, title: str, sentences: list[dict[str, Any]]) -> str:
    lines = [f"# {title}", ""]
    for sentence in sentences:
        lines.extend([
            f"[{format_timestamp(sentence['begin_time'])} – "
            f"{format_timestamp(sentence['end_time'])}] {sentence['text']}",
            "",
        ])
    content = "\n".join(lines).rstrip() + "\n"
    path.write_text(content, encoding="utf-8")
    return content


def main() -> int:
    for directory in (LOG_DIR, DIAG_DIR, OUTPUT_ROOT, PRIVATE_DIR, INBOX_DIR):
        directory.mkdir(parents=True, exist_ok=True)

    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    log_path = LOG_DIR / f"run-{stamp}.log"
    report_path = DIAG_DIR / f"run-report-{stamp}.json"

    def log(message: str) -> None:
        with log_path.open("a", encoding="utf-8") as handle:
            handle.write(f"{datetime.now(timezone.utc).isoformat()} {sanitize_text(message)}\n")

    report: dict[str, Any] = {
        "package": "V004-R2", "model": MODEL,
        "startedAt": datetime.now(timezone.utc).isoformat(), "status": "running"
    }

    try:
        config = configure_secrets()
        source = choose_one("最近的课程视频", {".mp4"})
        materials = choose_materials()
        course = input(f"课程名称 [{materials[0].stem[:40]}]: ").strip() or materials[0].stem[:40]
        lesson = input(f"课次名称 [{datetime.now().strftime('%Y-%m-%d')}]: ").strip() or datetime.now().strftime("%Y-%m-%d")
        raw_chunk = input("每段分钟数 [45]: ").strip()
        chunk_minutes = int(raw_chunk) if raw_chunk else 45

        lesson_id = safe_slug(f"{course}-{lesson}")
        output_dir = OUTPUT_ROOT / lesson_id
        maybe_import_previous_output(lesson_id, output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)
        chunks_dir = output_dir / "chunks"
        chunks_dir.mkdir(exist_ok=True)
        temp_dir = PRIVATE_DIR / lesson_id
        temp_dir.mkdir(parents=True, exist_ok=True)

        material_content, ocr_needed = write_materials(output_dir, materials)
        duration = ffprobe_duration(source, log)
        chunk_seconds = chunk_minutes * 60
        chunk_count = math.ceil(duration / chunk_seconds)
        print(f"\n→ 全课 {duration / 3600:.2f} 小时，共 {chunk_count} 段")

        endpoint = config.get("R2_ENDPOINT") or DEFAULT_R2_ENDPOINT
        bucket = config.get("R2_BUCKET") or DEFAULT_R2_BUCKET
        s3 = boto3.client(
            "s3", endpoint_url=endpoint,
            aws_access_key_id=config["R2_ACCESS_KEY_ID"],
            aws_secret_access_key=config["R2_SECRET_ACCESS_KEY"],
            region_name="auto", config=Config(signature_version="s3v4"),
        )

        checkpoints = []
        for zero in range(chunk_count):
            index = zero + 1
            checkpoint_path = chunks_dir / f"chunk-{index:03d}.json"
            task_path = chunks_dir / f"chunk-{index:03d}.task.json"
            start = zero * chunk_seconds
            length = min(chunk_seconds, duration - start)
            if checkpoint_path.exists():
                print(f"   ✓ 第 {index}/{chunk_count} 段已有结果，跳过")
                checkpoints.append(json.loads(checkpoint_path.read_text(encoding="utf-8")))
                continue
            print(f"→ 处理第 {index}/{chunk_count} 段")
            checkpoints.append(transcribe_chunk(
                source, index, start, length, checkpoint_path, task_path,
                temp_dir, config, s3, bucket, log
            ))
            print(f"   ✓ 第 {index} 段完成")

        sentences = []
        speech_ms = 0
        for checkpoint in sorted(checkpoints, key=lambda x: x["chunkIndex"]):
            sentences.extend(checkpoint["sentences"])
            speech_ms += int(checkpoint.get("speechDurationMilliseconds") or 0)
        sentences.sort(key=lambda x: (x["begin_time"], x["end_time"]))
        transcript = write_transcript(
            output_dir / "raw-transcript.md",
            f"{course} · {lesson} · 原始课堂转录",
            sentences,
        )
        lesson_input = (
            f"# {course} · {lesson} · AI 课程加工输入\n\n"
            "## 使用说明\n\n"
            "以下包含教师课件文字与 ASR 原始转录。课件用于确认规范术语、章节结构和材料内容；"
            "转录用于确认老师实际讲授、展开论证、案例和课堂边界。课件出现但转录无对应内容时，"
            "不得直接认定老师已经讲授。ASR 中可能存在同音错字，应结合课件语境理解，不要机械逐字照抄。\n\n"
            + material_content
            + "\n\n"
            + transcript
        )
        (output_dir / "lesson-input.md").write_text(lesson_input, encoding="utf-8")

        cost = speech_ms / 3_600_000 * PRICE_PER_HOUR_CNY
        summary = {
            "package": "V004-R2",
            "courseName": course,
            "lessonName": lesson,
            "videoDurationSeconds": round(duration, 3),
            "chunkCount": chunk_count,
            "sentenceCount": len(sentences),
            "characterCount": sum(len(x["text"]) for x in sentences),
            "ocrNeededCount": len(ocr_needed),
            "speechDurationMilliseconds": speech_ms,
            "estimatedCostCnyBeforeFreeQuota": round(cost, 6),
            "mechanicalCorrectionApplied": False,
        }
        (output_dir / "run-summary.json").write_text(
            json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
        )
        report.update({
            "status": "success",
            "finishedAt": datetime.now(timezone.utc).isoformat(),
            "outputDirectory": str(output_dir.relative_to(ROOT)),
            "summary": summary,
        })
        print("\n✓ V004 R2 完整课程输入已生成")
        print(f"句子：{len(sentences)}")
        print(f"待 OCR 页：{len(ocr_needed)}")
        print(f"免费额度前估算：¥{cost:.4f}")
        print(f"打开：{output_dir / 'lesson-input.md'}")
        return 0
    except Exception as exc:
        report.update({
            "status": "failed",
            "finishedAt": datetime.now(timezone.utc).isoformat(),
            "error": sanitize_text(str(exc)),
        })
        log("FAILED " + str(exc))
        print("\n✗ V004 R2 运行失败")
        print("原因：" + sanitize_text(str(exc)))
        print("已完成结果和已提交 task 都会保留，直接重跑即可继续。")
        return 1
    finally:
        report_path.write_text(
            json.dumps(sanitize_json(report), ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )


if __name__ == "__main__":
    raise SystemExit(main())
